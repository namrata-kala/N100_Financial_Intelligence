from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "N100 Financial Intelligence"
    subtitle.text = "Capstone Project Submission\nAutomated Analytics for India's Top 92 Companies"

    # Slide 2: Problem Statement
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Financial analysts spend hundreds of hours manually compiling data."
    p = tf.add_paragraph()
    p.text = "Our solution automates this entire pipeline."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "We built a scalable data platform to ingest, normalize, and analyze the top 92 Indian companies, generating real-time insights and automated PDF tearsheets."
    p.level = 1

    # Slide 3: Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "ETL Pipeline: Python (pandas) -> SQLite"
    p = tf.add_paragraph()
    p.text = "Machine Learning: scikit-learn (KMeans Clustering)"
    p = tf.add_paragraph()
    p.text = "Backend API: FastAPI (16 endpoints)"
    p = tf.add_paragraph()
    p.text = "Reporting: ReportLab (Automated PDFs)"
    p = tf.add_paragraph()
    p.text = "Testing: Pytest (74 robust unit/integration tests)"

    # Slide 4: Key Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Features & Deliverables"
    tf = body_shape.text_frame
    tf.text = "Over 100+ Automated Reports Generated (Company, Sector, Portfolio)"
    p = tf.add_paragraph()
    p.text = "Unsupervised Machine Learning for Company Archetype Clustering"
    p = tf.add_paragraph()
    p.text = "Dynamic Data Screener API"
    p = tf.add_paragraph()
    p.text = "Heuristic Rule Engine for Auto-Generating Pros & Cons"

    # Slide 5: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "All 20 Acceptance Gates Passed. Ready for Deployment."

    # Save
    os.makedirs("Final_Submission/PPT / Slides", exist_ok=True)
    prs.save("Final_Submission/PPT / Slides/N100_Financial_Intelligence_Presentation.pptx")
    print("Presentation saved!")

if __name__ == "__main__":
    create_presentation()
